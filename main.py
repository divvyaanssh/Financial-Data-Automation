import yfinance as yf
import pandas as pd
import os

tickers = ['AAPL', 'MSFT', 'TSLA']
data = []

for ticker in tickers:
    stock = yf.Ticker(ticker)
    info = stock.info
    pe = info.get("trailingPE", "N/A")
    roe = info.get("returnOnEquity", "N/A")
    growth = info.get("revenueGrowth", "N/A")
    data.append([ticker, pe, roe, growth])

df = pd.DataFrame(data, columns=['Ticker', 'P/E Ratio', 'ROE', 'Revenue Growth'])
df.to_excel('data/financial_data.xlsx', index=False)
from pptx import Presentation
from pptx.util import Inches

# Create output folder if it doesn't exist
os.makedirs('output', exist_ok=True)

# Create PowerPoint Presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide

# Add title
title_shape = slide.shapes.title
if not title_shape:
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
    title_frame = title_shape.text_frame
    title_frame.text = "Financial KPIs"
else:
    title_shape.text = "Financial KPIs"

# Calculate rows/columns
rows, cols = df.shape
table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8 + rows * 0.3)).table

# Header
for col in range(cols):
    table.cell(0, col).text = df.columns[col]

# Data
for row in range(rows):
    for col in range(cols):
        table.cell(row + 1, col).text = str(df.iloc[row, col])

# Save PPT
ppt_path = 'output/financial_summary.pptx'
prs.save(ppt_path)
print(f"✅ PowerPoint saved at {ppt_path}")

print("✅ Excel file created at data/financial_data.xlsx")
